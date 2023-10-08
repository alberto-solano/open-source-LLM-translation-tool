from pptx import Presentation
from utils import language_detection


def preprocess_text(doc_filepath, language):
    """Function to preprocess text from a .pptx file.

    Parameters
    ----------
    doc_filepath : str
        String with complete path to the document we aim to translate
    language : str | None
        Language of the original document if provided, if not it should be
        ``None'' and it will be inferred.

    Returns
    -------
    language : str
        Language of the document
    document : docx.document.Document
        Document object from .docx file containing all the text
    """
    # opens .pptx file
    document = Presentation(doc_filepath)

    # creates a list of non-empty paragraphs
    paragraphs_list = [
        shape.text_frame.text
        for slide in document.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]

    # infer language if it has not been specified
    if language is None:
        print("Detecting the language of the document...")
        language = language_detection(paragraphs_list)

    return language, document
